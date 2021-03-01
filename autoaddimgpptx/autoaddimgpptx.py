#python 3.7
#python-pptx 0.6.18
from pptx import Presentation
from pptx.util import Cm, Pt
from PIL import Image
import subprocess

# モジュールのインポート
import os, tkinter, tkinter.filedialog, tkinter.messagebox

# ファイル選択ダイアログの表示
root = tkinter.Tk()
root.withdraw()
fTyp = [("","*.jpg"), ("", "*.png")]
iDir = os.path.abspath(os.path.dirname(__file__))
#tkinter.messagebox.showinfo('○×プログラム','処理ファイルを選択してください！')

# ここの1行を変更 askopenfilename → askopenfilenames
file = tkinter.filedialog.askopenfilenames(filetypes = fTyp,initialdir = iDir)

# 選択ファイルリスト作成
#list = list(file)
#tkinter.messagebox.showinfo('○×プログラム',list)

prs = Presentation("template.pptx")
max_height = Cm(19.05)
max_width = Cm(33.867)
for i in range(len(file)):
    layout = prs.slide_layouts[6] # 6 白紙
    slide = prs.slides.add_slide(layout)
    img = Image.open(file[i])
    if img.width / img.height > 16 / 9:
        pic = slide.shapes.add_picture(file[i], 0, 0, width = max_width)
    else:
        pic = slide.shapes.add_picture(file[i], 0, 0, height = max_height)

prs.save("a.pptx")
subprocess.Popen([r"C:\Program Files\Microsoft Office\root\Office16\POWERPNT.EXE", r"a.pptx"])
